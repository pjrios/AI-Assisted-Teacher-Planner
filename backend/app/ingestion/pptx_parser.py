from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from ..schemas import YearlyPlan
from .text_parser import parse_yearly_plan_from_lines


def parse_yearly_plan_pptx(path: Path) -> YearlyPlan:
    presentation = Presentation(path)
    lines: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        lines.append(paragraph.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text
                        if text:
                            lines.extend(text.splitlines())
    return parse_yearly_plan_from_lines(lines)
